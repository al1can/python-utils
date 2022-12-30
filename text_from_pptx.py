import collections
import collections.abc
from pptx import Presentation
import glob
import os

def get_text_from_pptx():
    texts = ''
    for file in glob.glob("data/*.pptx"):
        prs = Presentation(file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    #text = ''.join(e for e in shape.text if e.isalnum() or e==' ')
                    escapes = ''.join([chr(char) for char in range(1, 32)])
                    text = shape.text.translate(str.maketrans('', '', escapes))
                    texts += text+'. '
    return texts
    